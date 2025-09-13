import os
import streamlit as st
import pandas as pd
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from pptx.dml.color import RGBColor
import math
import json
import re
# -----------------------------
# Initialize OpenAI client
# -----------------------------
api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)

# -----------------------------
# Load data
# -----------------------------
tasks_df = pd.read_csv("tasks.csv")
team_df = pd.read_csv("team.csv")
# Total project tasks
total_tasks = len(tasks_df)

# Completed tasks (past and current sprints)
completed_tasks = len(tasks_df[tasks_df['status'].isin(['done','completed'])])

# Current sprint number (assuming latest sprint is max sprint number)
current_sprint = tasks_df['sprint'].max()

# Tasks in current sprint
current_sprint_tasks_df = tasks_df[tasks_df['sprint'] == current_sprint]
sprint_tasks = len(current_sprint_tasks_df)
sprint_completed = len(current_sprint_tasks_df[current_sprint_tasks_df['status'].isin(['done','completed'])])

# Risk detection (example: if any task in current sprint has high priority and not completed)
risk_present = len(current_sprint_tasks_df[
    (current_sprint_tasks_df['priority'].str.lower() == 'high') & 
    (~current_sprint_tasks_df['status'].isin(['done','completed']))
]) > 0

@st.cache_data
def load_data():
    tasks_df = pd.read_csv("tasks.csv")
    team_df = pd.read_csv("team.csv")
    merged = tasks_df.merge(team_df, on="owner", how="left")
    return merged

# -----------------------------
# Streamlit UI
# -----------------------------
st.set_page_config(page_title="📊PM Project Lens", layout="wide")
col1, col2 = st.columns([1, 6])  # adjust width ratio
with col1:
    st.image("pmprojectlens.png", width=160)  # your logo
with col2:
    st.title("PM Project Lens - AI Powered Project Manager Assistant")
    st.subheader("Weekly Status Report (WSR), Risk Analysis & Mitigation Plan")  # Additional text


prompt_input = st.text_area(
    "Enter your risk analysis prompt", 
    "Create Sprint 5 WSR and Identify tasks that are most likely to cause delays due to vague requirements or high priority."
)

# 🔹 Slider to control number of records
max_records = st.slider(
    "Select number of records to display:",
    min_value=1,
    max_value=len(tasks_df),
    value=min(5, len(tasks_df)),  # default 5 or less if fewer tasks
    step=1
)
def remove_slide(prs, slide):
    """Remove a slide from a presentation"""
    xml_slides = prs.slides._sldIdLst  # internal slide list
    slides = list(xml_slides)
    for i, s in enumerate(slides):
        if prs.slides[i] == slide:
            xml_slides.remove(s)
            break
            
def extract_sprint_number(prompt_text):
    """
    Try to extract sprint number from the user prompt.
    Matches: 'sprint 5', 'sprint #5', 'Sprint Number 5', 'sprint:5'
    Returns int or None.
    """
    if not prompt_text:
        return None
    patterns = [
        r"\bsprint\s*#?\s*(\d+)\b",
        r"\bsprint\s*number\s*#?\s*(\d+)\b",
        r"\biteration\s*#?\s*(\d+)\b"
    ]
    for pat in patterns:
        m = re.search(pat, prompt_text, re.IGNORECASE)
        if m:
            try:
                return int(m.group(1))
            except:
                pass
    return None

def calculate_project_overview(total_tasks, completed_tasks, sprint_tasks, sprint_completed, risk_present):
    # Project scope delivered %
    project_scope_delivered = (completed_tasks / total_tasks) * 100 if total_tasks else 0
    
    # Sprint % delivered
    sprint_delivered = (sprint_completed / sprint_tasks) * 100 if sprint_tasks else 0
    
    # Determine Project Status
    if completed_tasks == total_tasks and not risk_present:
        project_status = ("Good", RGBColor(0, 176, 80))  # Green
    elif completed_tasks == total_tasks and risk_present or completed_tasks != total_tasks and not risk_present:
        project_status = ("Average", RGBColor(255, 192, 0))  # Amber
    else:
        project_status = ("Bad", RGBColor(255, 0, 0))  # Red
    
    return project_scope_delivered, sprint_delivered, project_status


def add_project_overview_slide(prs, total_tasks, completed_tasks, sprint_tasks, sprint_completed, risk_present):
    # Calculate values
    project_scope, sprint_pct, status = calculate_project_overview(
        total_tasks, completed_tasks, sprint_tasks, sprint_completed, risk_present
    )
    status_text, status_color = status

    # Add slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    # Remove all placeholder shapes (like "Click to add title")
    for shape in slide.shapes:
        if shape.is_placeholder:
            sp = shape
            slide.shapes._spTree.remove(sp._element)
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Project Overview"
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True

    # Project Scope Delivered
    scope_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    scope_tf = scope_box.text_frame
    scope_tf.text = f"PROJECT SCOPE DELIVERED: {project_scope:.2f}%"

    # Sprint % Delivered
    sprint_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1))
    sprint_tf = sprint_box.text_frame
    sprint_tf.text = f"SPRINT % DELIVERED: {sprint_pct:.2f}%"

    # Project Status
    status_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(1))
    status_tf = status_box.text_frame
    p = status_tf.paragraphs[0]
    p.text = f"PROJECT STATUS: {status_text}"
    p.font.color.rgb = status_color
    p.font.bold = True

def create_wsr_pptx_from_df(results_df,
                            sprint_df=None,
                            sprint_number=None,
                            template_path="WSR_Framework.pptx",
                            title_text="Risks and Mitigation Plan",
                            max_rows_per_slide=12):
    """
    Create a PPTX with:
      - Slide 1: Sprint Overview (S.No, Task Details, Status)  [if sprint_df + sprint_number provided]
      - Slide 2+: Risks and Mitigation Plan (task_name, owner, risk_score, mitigation_plan)

    Parameters:
      results_df: DataFrame with columns: task_name, owner, risk_score, mitigation_plan
      sprint_df: DataFrame with columns containing task_name/description and status (optional)
      sprint_number: int or str for slide title (optional)
      template_path: path to template pptx (optional)
      title_text: title used for the Risks slides (kept for backward compatibility)
      max_rows_per_slide: how many table rows per slide
    Returns:
      bytes of PPTX file
    """
    # local imports so missing libs don't break on import-time
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from io import BytesIO
    import math
    import pandas as pd

    # Validate risk columns
    expected_cols = ["task_name", "owner", "risk_score", "mitigation_plan"]
    for c in expected_cols:
        if c not in results_df.columns:
            raise ValueError(f"Missing column in results_df: {c}")

    # Load template or blank presentation
    try:
        prs = Presentation(template_path) if template_path else Presentation()
    except Exception:
        prs = Presentation()

    # Helper: add a slide with a table using a blank-like layout (avoids placeholder overlap)
    def _add_table_slide(prs, df_chunk, slide_title, headers, col_width_ratios):
        # choose a layout with minimal placeholders (best-effort)
        layout = None
        for l in prs.slide_layouts:
            if len(l.placeholders) == 0:
                layout = l
                break
        if layout is None:
            layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        # Remove all placeholder shapes (like "Click to add title")
        for shape in slide.shapes:
            if shape.is_placeholder:
                sp = shape
                slide.shapes._spTree.remove(sp._element)

        # Add our own title textbox to avoid built-in placeholder conflicts
        left = Inches(0.5); top = Inches(0.2); width = Inches(9); height = Inches(0.6)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_tf.text = slide_title
        title_tf.paragraphs[0].font.size = Pt(22)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # compute slide width in inches (safe)
        try:
            slide_width_inches = prs.slide_width / 914400.0
        except Exception:
            slide_width_inches = 10.0

        table_width_inches = slide_width_inches - 0.8
        table_width = Inches(table_width_inches)

        rows = len(df_chunk) + 1
        cols = len(headers)
        left = Inches(0.4); top = Inches(1.0); height = Inches(0.5 + 0.28 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, table_width, height).table

        # set column widths by ratios
        for i, ratio in enumerate(col_width_ratios):
            try:
                table.columns[i].width = Inches(table_width_inches * ratio)
            except Exception:
                pass

        # header row
        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)

        # fill rows (expect df_chunk columns aligned to headers order)
        for r, (_, row) in enumerate(df_chunk.iterrows(), start=1):
            for ci, col in enumerate(df_chunk.columns[:cols]):
                table.cell(r, ci).text = str(row[col])
                table.cell(r, ci).text_frame.paragraphs[0].font.size = Pt(10)

    # --- Slide 1: Sprint Overview (optional) ---
    if sprint_df is not None and sprint_number is not None:
        sprint_df_local = sprint_df.copy()
        # Build "Task Details" using description if present otherwise task_name
        if "description" in sprint_df_local.columns:
            sprint_df_local["Task Details"] = sprint_df_local["description"].astype(str)
        else:
            sprint_df_local["Task Details"] = sprint_df_local.get("task_name", "").astype(str)
        sprint_df_local["Status"] = sprint_df_local.get("status", "Unknown").astype(str)

        sprint_df_local = sprint_df_local.reset_index(drop=True)
        sprint_df_local.index = sprint_df_local.index + 1
        sprint_table = pd.DataFrame({
            "S.No": sprint_df_local.index.astype(str),
            "Task Details": sprint_df_local["Task Details"].astype(str),
            "Status": sprint_df_local["Status"].astype(str)
        })

        _add_table_slide(
            prs,
            sprint_table,
            slide_title=f"Sprint {sprint_number} - Overview",
            headers=["S.No", "Task Details", "Status"],
            col_width_ratios=[0.08, 0.78, 0.14]
        )

    # --- Slide 2+: Risks and Mitigation Plan ---
    total_rows = len(results_df)
    # Ensure at least one slide even if empty
    n_slides = math.ceil(total_rows / max_rows_per_slide) if total_rows > 0 else 1
    for i in range(n_slides):
        start = i * max_rows_per_slide
        end = start + max_rows_per_slide
        chunk = results_df.iloc[start:end][["task_name", "owner", "risk_score", "mitigation_plan"]].copy()
        chunk.columns = ["Task", "Owner", "Risk Score", "Mitigation Plan"]

        _add_table_slide(
            prs,
            chunk,
            slide_title=title_text if i == 0 else title_text,  # keep same title for each risk slide
            headers=["Task", "Owner", "Risk Score", "Mitigation Plan"],
            col_width_ratios=[0.25, 0.15, 0.12, 0.48]
        )
    
    # Completed tasks (status == Closed)
    completed_tasks = len(tasks_df[tasks_df['status'] == 'Closed'])

    # Current sprint tasks
    current_sprint_tasks_df = tasks_df[tasks_df['sprint'] == sprint_number]
    sprint_tasks = len(current_sprint_tasks_df)
    sprint_completed = len(current_sprint_tasks_df[current_sprint_tasks_df['status'] == 'Closed'])

    # Risk present
    risk_present = len(current_sprint_tasks_df[
        (current_sprint_tasks_df['priority'].str.lower() == 'high') & 
        (current_sprint_tasks_df['status'].isin(['Open','InProgress']))
    ]) > 0

    add_project_overview_slide(prs, total_tasks, completed_tasks, sprint_tasks, sprint_completed, risk_present)
    remove_slide(prs, prs.slides[1])
    # Save to bytes
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()



if st.button("Analyze Risks"):
    with st.spinner("Analyzing tasks..."):
        if not prompt_input.strip():
            st.warning("Please enter a prompt for risk analysis.")
        else:
            # -----------------------------
            # Prepare tasks data for LLM
            # -----------------------------
            tasks_json = tasks_df.to_dict(orient="records")
            team_json = team_df.to_dict(orient="records")

            system_message = f"""
            You are a project management assistant. 
            Given a list of tasks and team info, analyze the risk for each task based on the following criteria:
            - Vague requirements
            - High priority
            - Capacity of assigned owner

            Data provided:
            Tasks: {json.dumps(tasks_json)}
            Team: {json.dumps(team_json)}
            Filter only Top {max_records} records
            Your output must be JSON array with fields:
            - task_name
            - owner
            - risk_score (High, Medium, Low)
            - mitigation_plan (suggested actions to reduce risk)
            """
            
            try:
                response = client.chat.completions.create(
                    model="gpt-4.1-mini",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt_input}
                    ],
                    temperature=0.3,
                )
                
                output_text = response.choices[0].message.content
                #st.subheader("Risk Analysis Output (JSON)")
                #st.code(output_text, language="json")
                
                # Try parsing JSON output into DataFrame
                try:
                    output_json = json.loads(output_text)
                    output_df = pd.DataFrame(output_json)
                    
                     # 🔹 Apply record count filter
                    filtered_results = output_df.head(max_records)
                    st.subheader("Risk Analysis Table")
                    #st.dataframe(output_df)
                    filtered_results.insert(0, "S.No", range(1, 1 + len(filtered_results)))
                    st.dataframe(filtered_results, hide_index=True, use_container_width=True)
                    final_df = filtered_results.copy()
                    # Optionally allow CSV download
                    csv = output_df.to_csv(index=False).encode('utf-8')
                    st.download_button(
                        "Download CSV", csv, "risk_analysis.csv", "text/csv", key='download-csv'
                    )
                    for col in ["task_name", "owner", "risk_score", "mitigation_plan"]:
                        if col not in final_df.columns:
                            final_df[col] = ""

                    # Create PPTX bytes
                    pptx_bytes = None
                    try:
                        pptx_bytes = create_wsr_pptx_from_df(final_df, template_path="WSR_Framework.pptx",
                                                             title_text="Risks and Mitigation Plan",
                                                             max_rows_per_slide=12)
                    except Exception as e:
                        st.error(f"Failed to generate PPTX from template: {e}")
                        # Try without template
                        pptx_bytes = create_wsr_pptx_from_df(final_df, template_path=None,
                                                             title_text="Risks and Mitigation Plan",
                                                             max_rows_per_slide=12)
                    data_df = load_data()
                    sprint_num = extract_sprint_number(prompt_input)
                    sprint_df = data_df[data_df["sprint"] == sprint_num][["task_name","description","status"]].copy()
                    
                    if sprint_num is None:
                        sprint_num = st.sidebar.number_input("Sprint number (fallback)", min_value=1, value=1)

                    pptx_bytes = create_wsr_pptx_from_df(
                        results_df=final_df[["task_name","owner","risk_score","mitigation_plan"]],
                        sprint_df=sprint_df,
                        sprint_number=sprint_num,
                        template_path="WSR_Framework.pptx"
                    )
                    
                    # 4) download button
                    st.download_button(
                        f"⬇ Download WSR for Sprint {sprint_num}",
                        data=pptx_bytes,
                        file_name=f"WSR_Sprint_{sprint_num}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"Failed to parse JSON output: {e}")

            except Exception as e:
                st.error(f"Error calling OpenAI API: {e}")
